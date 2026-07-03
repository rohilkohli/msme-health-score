"""
Generate the IDBI Innovate 2026 submission deck from the provided template.
Run: python generate_deck.py
Output: MSME_Financial_Health_Score_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

TEMPLATE_PATH = os.path.expanduser("~/Downloads/Prototype Submission Deck _ IDBI Innovate.pptx")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "MSME_Financial_Health_Score_Deck.pptx")

BLUE = RGBColor(0x1E, 0x3A, 0x5F)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_text(shape, text, font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment


def add_bullet_text(text_frame, bullets, font_size=12, color=DARK):
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def fill_slide_1(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if "team name" in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "Team Name: FinScore AI"
                p.font.size = Pt(16)
                p.font.bold = True
            elif "team leader" in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "Team Leader: Rohil Kohli"
                p.font.size = Pt(16)
            elif "problem statement" in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "Problem Statement 3: Financial Health Score"
                p.font.size = Pt(16)
                p.font.color.rgb = BLUE


def fill_slide_2(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if "brief" in text or shape.text_frame.text.strip() == "":
                if "brief" in text:
                    continue
                add_bullet_text(shape.text_frame, [
                    "An AI/ML-driven MSME Financial Health Card platform that aggregates alternate data from GST, UPI, EPFO, and Account Aggregator ecosystems.",
                    "",
                    "It computes a multidimensional financial health score (0-1000) across 5 key dimensions: Revenue Stability, Cash Flow Health, Compliance Score, Growth Trajectory, and Repayment Capacity.",
                    "",
                    "The platform visualizes strengths and risks through an intuitive Health Card interface, integrates with India Stack (ULI/OCEN/AA), and enables near real-time credit assessment for credit-invisible MSMEs.",
                    "",
                    "This directly addresses the challenge of 63M+ Indian MSMEs being excluded from formal credit due to lack of traditional financial documentation, while leveraging India's world-class digital public infrastructure."
                ], font_size=11)


def fill_slide_3(slide):
    content = {
        "how different": "Unlike traditional credit scoring (CIBIL/bureau-dependent), our solution works for New-to-Credit (NTC) and New-to-Bank (NTB) MSMEs by leveraging alternate data. No existing solution provides a unified 5-dimension scoring framework that integrates all of GST + UPI + EPFO + AA in real-time with explainable AI.",
        "how will it": "By creating a consent-based, privacy-first data aggregation pipeline that pulls from India Stack APIs, computes scores using ML models trained on MSME behavioral patterns, and delivers actionable insights — both to the MSME (improvement roadmap) and to lenders (standardized risk assessment via OCEN/ULI).",
        "usp": "1) First unified 5-dimension scoring for NTC MSMEs\n2) Real-time scoring (not batch)\n3) Explainable AI with dimension-wise breakdown\n4) Full India Stack integration (AA/ULI/OCEN)\n5) Works with ZERO traditional documents\n6) Self-improving ML model with feedback loops"
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if "different" in text:
                continue
            elif "solve" in text:
                continue
            elif "usp" in text:
                continue


def fill_slide_4(slide):
    features = [
        "• Consent-based Account Aggregator data fetching (RBI compliant)",
        "• Real-time GST filing analysis — revenue trends, filing consistency, seasonal patterns",
        "• UPI transaction intelligence — cash flow patterns, customer diversity, payment regularity",
        "• EPFO compliance tracking — employee contribution consistency, workforce stability",
        "• 5-Dimension Financial Health Score (0-1000) with category classification",
        "• Interactive Health Card with radar visualization and trend analysis",
        "• AI-powered credit readiness assessment with eligible loan products",
        "• Personalized improvement recommendations with priority ranking",
        "• Score history tracking with temporal trend analysis",
        "• ULI/OCEN-ready API for seamless lender integration",
        "• Portfolio-level analytics dashboard for bank officers",
        "• Risk pattern detection using ML anomaly detection",
        "• Industry benchmark comparison for contextual scoring",
        "• Multi-language support (English + Hindi)"
    ]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if "list of features" in text:
                continue
            elif shape.text_frame.text.strip() == "" or "feature" not in text:
                add_bullet_text(shape.text_frame, features, font_size=10)
                break


def fill_slide_8(slide):
    tech = [
        "BACKEND:",
        "• Python 3.11 + FastAPI (high-performance async API framework)",
        "• SQLAlchemy + SQLite/PostgreSQL (database ORM)",
        "• scikit-learn + pandas + numpy (ML & data processing)",
        "• python-jose (JWT authentication)",
        "• Pydantic (data validation & serialization)",
        "",
        "FRONTEND:",
        "• React 18 + Vite (modern SPA framework)",
        "• TailwindCSS (utility-first responsive styling)",
        "• Recharts (data visualization — radar, line, pie charts)",
        "• Framer Motion (smooth animations)",
        "• Axios (API communication)",
        "",
        "INFRASTRUCTURE:",
        "• Docker + Docker Compose (containerization)",
        "• GitHub Actions (CI/CD pipeline)",
        "• Vercel (frontend deployment)",
        "• Railway/Render (backend deployment)",
        "",
        "ECOSYSTEM INTEGRATION:",
        "• Account Aggregator APIs (Setu/Finvu sandbox)",
        "• GST APIs (government portal integration)",
        "• OCEN protocol (standardized credit APIs)",
        "• ULI framework (unified lending interface)"
    ]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if "technologies" in text:
                continue
            else:
                add_bullet_text(shape.text_frame, tech, font_size=9)
                break


def fill_slide_12(slide):
    future = [
        "FUTURE DEVELOPMENT:",
        "",
        "• Integration with live Account Aggregator FIPs (Setu, Finvu, OneMoney)",
        "• Real-time GST API integration with government portal",
        "• Advanced ML models: Gradient Boosted Trees + Deep Learning ensemble",
        "• Supply chain graph analysis for network-based scoring",
        "• WhatsApp bot for score delivery and recommendations",
        "• Multi-language voice interface for Tier-2/3 city MSMEs",
        "• Blockchain-based score attestation for tamper-proof sharing",
        "• Integration with GeM (Government e-Marketplace) transaction data",
        "• Sector-specific scoring models (manufacturing, services, trading)",
        "• Continuous score monitoring with alert system",
        "• Partnership API for NBFCs and fintech lenders",
        "• Regulatory reporting dashboard for RBI compliance"
    ]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if "additional" in text or "future" in text:
                continue
            else:
                add_bullet_text(shape.text_frame, future, font_size=10)
                break


def fill_slide_13(slide):
    links = [
        "GitHub Repository: https://github.com/YOUR_USERNAME/msme-health-score",
        "",
        "Demo Video Link: [3-minute demo video URL]",
        "",
        "Live Application: [Deployment URL]"
    ]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if "provide links" in text or "github" in text:
                continue
            else:
                add_bullet_text(shape.text_frame, links, font_size=12)
                break


def main():
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)

    if len(slides) >= 1:
        fill_slide_1(slides[0])
    if len(slides) >= 2:
        fill_slide_2(slides[1])
    if len(slides) >= 3:
        fill_slide_3(slides[2])
    if len(slides) >= 4:
        fill_slide_4(slides[3])
    if len(slides) >= 8:
        fill_slide_8(slides[7])
    if len(slides) >= 12:
        fill_slide_12(slides[11])
    if len(slides) >= 13:
        fill_slide_13(slides[12])

    prs.save(OUTPUT_PATH)
    print(f"Deck saved to: {OUTPUT_PATH}")
    print("Note: Add architecture diagrams, wireframes, and screenshots manually or re-run after building the app.")


if __name__ == "__main__":
    main()
